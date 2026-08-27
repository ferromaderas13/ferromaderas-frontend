# -*- coding: utf-8 -*-
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

BASE = Path(
    r"C:\Users\geova\OneDrive\Desktop\10 semestre\Proyecto de graduación II"
    r"\Version editada PG2\SEGUNDA ENTREGA - CAPITULO IV"
)
IMG = BASE / "figura-01-arquitectura-ppt.png"
OUT = BASE / "Presentacion-Ferromaderas-chilera.pptx"

TEAL = RGBColor(0x0F, 0x3D, 0x4C)
TEAL_DEEP = RGBColor(0x0B, 0x2E, 0x3A)
TEAL_MID = RGBColor(0x1A, 0x6B, 0x78)
GREEN = RGBColor(0x3D, 0xDC, 0x97)
GREY = RGBColor(0x6B, 0x72, 0x80)
GREY_SOFT = RGBColor(0xEE, 0xF1, 0xF4)
PAGE = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x11, 0x27, 0x33)
MINT = RGBColor(0xB8, 0xD4, 0xDC)

SERIF = "Georgia"
SANS = "Calibri"


def no_line(shape):
    shape.line.fill.background()


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    no_line(shape)


def round_rect(slide, l, t, w, h, color, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(sh, color)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def run(p, text, size, bold=False, color=INK, font=SANS, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def page_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(bg, color)
    return bg


def title_block(slide, title, subtitle, y0=0.28):
    tf = tb(slide, Inches(0.7), Inches(y0), Inches(12), Inches(0.58))
    p = tf.paragraphs[0]
    run(p, title, 30, True, TEAL, SERIF)
    tf = tb(slide, Inches(0.7), Inches(y0 + 0.58), Inches(12), Inches(0.4))
    p = tf.paragraphs[0]
    run(p, subtitle, 14, False, GREY, SANS)


def arrow(slide, l, t):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, Inches(0.42), Inches(0.22))
    fill(sh, TEAL_MID)
    return sh


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ========== 1 PORTADA ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, WHITE)

round_rect(s, Inches(0.85), Inches(0.7), Inches(3.35), Inches(0.38), GREY_SOFT, 0.5)
tf = tb(s, Inches(0.85), Inches(0.74), Inches(3.35), Inches(0.34))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "Proyecto de graduación II", 12, True, TEAL, SANS)

tf = tb(s, Inches(0.85), Inches(1.35), Inches(11.4), Inches(1.9))
p = tf.paragraphs[0]
run(p, "Plataforma web\npara Ferromaderas", 44, True, TEAL, SERIF)

tf = tb(s, Inches(0.85), Inches(3.45), Inches(8.6), Inches(1.7))
p = tf.paragraphs[0]
run(
    p,
    "Desarrollo de una plataforma web responsive con cotización en línea, asistente conversacional inteligente y recomendación de productos para la consulta comercial y el seguimiento de clientes en Ferromaderas, Amatitlán.",
    16,
    False,
    GREY,
    SANS,
)

round_rect(s, Inches(0.85), Inches(5.45), Inches(11.6), Inches(1.45), TEAL, 0.08)
tf = tb(s, Inches(1.15), Inches(5.62), Inches(11.0), Inches(1.15))
p = tf.paragraphs[0]
run(p, "Geovany Emmanuel González Díaz", 18, True, WHITE, SANS)
p2 = tf.add_paragraph()
run(
    p2,
    "Universidad Mariano Gálvez de Guatemala  ·  Ingeniería en Sistemas  ·  Ferromaderas, Amatitlán",
    14,
    False,
    MINT,
    SANS,
)

# ========== 2 OBJETIVO ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, WHITE)
title_block(
    s,
    "¿Cuál es el objetivo?",
    "Una sola plataforma web, en celular y computadora, para cotizar, orientar al cliente y dar seguimiento.",
)

items = [
    ("01", "Cotización en línea", "El cliente arma y comparte su cotización sin pasar por WhatsApp suelto."),
    ("02", "Asistente y recomendaciones", "Consulta productos y recibe sugerencias desde el catálogo real."),
    ("03", "Seguimiento interno", "El personal ve estados, responsables y el hilo de cada solicitud."),
    ("04", "Tablero e indicadores", "Gerencia consulta carga, conversión, reportes y analítica del mismo negocio."),
]
for i, (num, title, desc) in enumerate(items):
    y = Inches(1.62) + Inches(i * 1.25)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.58), Inches(0.58))
    fill(circ, GREY_SOFT)
    tf = tb(s, Inches(0.7), y + Inches(0.12), Inches(0.58), Inches(0.38))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, num, 12, True, TEAL_MID, SANS)
    tf = tb(s, Inches(1.5), y - Inches(0.02), Inches(5.3), Inches(1.1))
    p = tf.paragraphs[0]
    run(p, title, 16, True, TEAL, SANS)
    p2 = tf.add_paragraph()
    run(p2, desc, 13, False, GREY, SANS)

round_rect(s, Inches(7.15), Inches(1.62), Inches(5.45), Inches(5.15), TEAL_DEEP, 0.07)
tf = tb(s, Inches(7.45), Inches(1.88), Inches(4.9), Inches(0.35))
p = tf.paragraphs[0]
run(p, "Objetivo general", 12, True, GREEN, SANS)

tf = tb(s, Inches(7.45), Inches(2.32), Inches(4.9), Inches(3.7))
p = tf.paragraphs[0]
run(
    p,
    "Desarrollar e implementar una plataforma web responsive con cotización en línea, asistente conversacional inteligente, recomendación de productos, seguimiento comercial estructurado y tablero de indicadores, orientada a optimizar la consulta de productos, el control de cotizaciones y la atención al cliente en Ferromaderas, Amatitlán.",
    16,
    False,
    WHITE,
    SERIF,
)
for para in tf.paragraphs:
    para.line_spacing = 1.25

tf = tb(s, Inches(7.45), Inches(6.18), Inches(4.9), Inches(0.4))
p = tf.paragraphs[0]
run(p, "Capítulo I · mismo texto de la tesis", 11, False, RGBColor(0x8A, 0xB4, 0xBE), SANS, True)

# ========== 3 ARQUITECTURA (layout tipo referencia) ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, PAGE)
title_block(
    s,
    "Arquitectura del sistema",
    "Tres capas: lo que ve el usuario, la API de negocio y los datos. Abajo, lo que sostiene el despliegue y la medición.",
)

stack = [
    ("Frontend", "Angular", "Catálogo, cotización, asistente, recomendaciones y panel interno. Responsive en celular y computadora."),
    ("Backend", "NestJS + Express", "API REST: 2FA, cotizaciones, inventario, alertas, bitácora y reportes operativos."),
    ("Datos", "PostgreSQL", "Productos, cotizaciones, usuarios y seguimiento. Inventario alineado con Di-Chara."),
]
box_w = Inches(3.55)
box_h = Inches(3.05)
gap = Inches(0.55)
start_x = Inches(0.7)
y_box = Inches(1.55)
for i, (kicker, title, desc) in enumerate(stack):
    x = start_x + i * (box_w + gap)
    round_rect(s, x, y_box, box_w, box_h, WHITE, 0.06)
    circ = s.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(1.38), y_box + Inches(0.28), Inches(0.78), Inches(0.78)
    )
    fill(circ, RGBColor(0xE6, 0xF2, 0xF4))
    tf = tb(s, x + Inches(1.38), y_box + Inches(0.46), Inches(0.78), Inches(0.48))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, ("0%d" % (i + 1)), 14, True, TEAL_MID, SANS)
    tf = tb(s, x + Inches(0.22), y_box + Inches(1.18), Inches(3.1), Inches(1.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, kicker, 11, True, TEAL_MID, SANS)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run(p2, title, 18, True, TEAL, SERIF)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run(p3, desc, 12, False, GREY, SANS)
    if i < 2:
        arrow(s, x + box_w + Inches(0.06), y_box + Inches(1.4))

support = [
    ("GitHub + CI/CD", "Despliegue continuo: frontend en Vercel y API en Railway."),
    ("Analítica web", "GTM y GA4: visitas, búsqueda, chatbot y eventos de cotización."),
    ("Power BI", "Tablero gerencial sobre los mismos datos: carga, conversión e indicadores."),
]
sy = Inches(5.0)
sw = Inches(3.85)
sg = Inches(0.25)
for i, (title, desc) in enumerate(support):
    x = Inches(0.7) + i * (sw + sg)
    round_rect(s, x, sy, sw, Inches(1.85), WHITE, 0.08)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), sy + Inches(0.38), Inches(0.48), Inches(0.48))
    fill(circ, GREY_SOFT)
    tf = tb(s, x + Inches(0.22), sy + Inches(0.48), Inches(0.48), Inches(0.32))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, ("0%d" % (i + 1)), 10, True, TEAL_MID, SANS)
    tf = tb(s, x + Inches(0.85), sy + Inches(0.32), Inches(2.8), Inches(1.3))
    p = tf.paragraphs[0]
    run(p, title, 15, True, TEAL, SANS)
    p2 = tf.add_paragraph()
    run(p2, desc, 12, False, GREY, SANS)

# ========== 4 DESPLIEGUE (Figura 1 tesis) ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, WHITE)
title_block(
    s,
    "¿Cómo está desplegada?",
    "Figura 1. Diagrama de arquitectura de despliegue de la plataforma web de Ferromaderas.",
)
if IMG.exists():
    pic = s.shapes.add_picture(str(IMG), Inches(0.55), Inches(1.22), width=Inches(12.2))
    max_h = Inches(5.9)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_h
        pic.left = int((W - pic.width) / 2)
else:
    tf = tb(s, Inches(1), Inches(3), Inches(10), Inches(1))
    run(tf.paragraphs[0], "No se encontró la Figura 1. Pegá el PNG a mano.", 16, False, GREY, SANS)

# ========== 5 RESULTADOS ESPERADOS ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, WHITE)
title_block(
    s,
    "¿Qué resultados se esperan?",
    "Que Ferromaderas deje de cotizar y dar seguimiento a pedazos, y lo haga sobre un solo sistema.",
)

results = [
    ("Consulta comercial", "El cliente encuentra producto, precio y stock sin depender de un mensaje suelto."),
    ("Cotización trazable", "Cada solicitud queda con código, estado y responsable. No se pierde en el chat."),
    ("Atención guiada", "El asistente y las recomendaciones usan el catálogo real, no una lista inventada."),
    ("Visión de gerencia", "Reportes operativos, analítica web y tablero Power BI sobre el mismo negocio."),
]
for i, (title, desc) in enumerate(results):
    y = Inches(1.62) + Inches(i * 1.22)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.58), Inches(0.58))
    fill(circ, GREY_SOFT)
    tf = tb(s, Inches(0.7), y + Inches(0.12), Inches(0.58), Inches(0.38))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "0%d" % (i + 1), 12, True, TEAL_MID, SANS)
    tf = tb(s, Inches(1.5), y - Inches(0.02), Inches(5.35), Inches(1.1))
    p = tf.paragraphs[0]
    run(p, title, 16, True, TEAL, SANS)
    p2 = tf.add_paragraph()
    run(p2, desc, 13, False, GREY, SANS)

round_rect(s, Inches(7.15), Inches(1.62), Inches(5.45), Inches(5.15), TEAL_DEEP, 0.07)
tf = tb(s, Inches(7.45), Inches(1.88), Inches(4.9), Inches(0.35))
p = tf.paragraphs[0]
run(p, "En la práctica", 12, True, GREEN, SANS)

tf = tb(s, Inches(7.45), Inches(2.35), Inches(4.9), Inches(3.5))
p = tf.paragraphs[0]
run(
    p,
    "Menos cotizaciones perdidas, menos transcripción a mano y una lectura clara de qué se cotiza, qué se cierra y qué hay que atender.",
    20,
    False,
    WHITE,
    SERIF,
)
for para in tf.paragraphs:
    para.line_spacing = 1.3

tf = tb(s, Inches(7.45), Inches(6.15), Inches(4.9), Inches(0.4))
p = tf.paragraphs[0]
run(p, "Eso es lo que se muestra en la demo.", 12, False, RGBColor(0x8A, 0xB4, 0xBE), SANS, True)

# ========== 6 ANALÍTICA Y POWER BI ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, PAGE)
title_block(
    s,
    "Analítica y tablero gerencial",
    "Tres lecturas del mismo negocio: lo que pasa en el sitio, lo operativo del día y lo que ve gerencia.",
)

cols = [
    (
        "Analítica web",
        "GA4 + GTM",
        "Visitas, búsqueda, ficha de producto, chatbot y eventos al generar o compartir una cotización.",
    ),
    (
        "Reportes internos",
        "Panel /admin/reportes",
        "Estados de cotización, productos más cotizados, vendedores y exportación a PDF.",
    ),
    (
        "Power BI",
        "Tablero gerencial",
        "Indicadores de carga, conversión y seguimiento. Acá entra la demo en vivo.",
    ),
]
cw = Inches(3.85)
for i, (kicker, title, desc) in enumerate(cols):
    x = Inches(0.7) + i * (cw + Inches(0.25))
    round_rect(s, x, Inches(1.55), cw, Inches(4.55), WHITE, 0.07)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.53), Inches(1.85), Inches(0.78), Inches(0.78))
    fill(circ, RGBColor(0xE6, 0xF2, 0xF4))
    tf = tb(s, x + Inches(1.53), Inches(2.04), Inches(0.78), Inches(0.45))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "0%d" % (i + 1), 14, True, TEAL_MID, SANS)
    tf = tb(s, x + Inches(0.28), Inches(2.85), Inches(3.28), Inches(2.9))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, kicker, 12, True, TEAL_MID, SANS)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run(p2, title, 18, True, TEAL, SERIF)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run(p3, desc, 13, False, GREY, SANS)

tf = tb(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.55))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "En la siguiente diapositiva se abre la demo. Power BI y el panel los mostrás vos en vivo.", 13, False, GREY, SANS)

# ========== 7 DEMO ==========
s = prs.slides.add_slide(prs.slide_layouts[6])
page_bg(s, WHITE)

round_rect(s, Inches(0.85), Inches(0.7), Inches(2.15), Inches(0.38), GREY_SOFT, 0.5)
tf = tb(s, Inches(0.85), Inches(0.74), Inches(2.15), Inches(0.34))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "En vivo", 12, True, TEAL, SANS)

tf = tb(s, Inches(0.85), Inches(1.3), Inches(11.4), Inches(1.5))
p = tf.paragraphs[0]
run(p, "Demo de la plataforma", 40, True, TEAL, SERIF)

tf = tb(s, Inches(0.85), Inches(2.9), Inches(11.4), Inches(0.55))
p = tf.paragraphs[0]
run(p, "https://ferromaderas-frontend.vercel.app/", 18, False, TEAL_MID, SANS)

beats = [
    ("01", "Cliente", "Catálogo, asistente, recomendaciones y cotización."),
    ("02", "Interno", "Seguimiento, alertas, 2FA y reportes del panel."),
    ("03", "Gerencia", "Indicadores y tablero Power BI."),
]
for i, (num, title, desc) in enumerate(beats):
    x = Inches(0.85) + i * Inches(4.0)
    round_rect(s, x, Inches(3.7), Inches(3.75), Inches(2.35), GREY_SOFT, 0.08)
    tf = tb(s, x + Inches(0.28), Inches(3.9), Inches(3.2), Inches(1.95))
    p = tf.paragraphs[0]
    run(p, num, 12, True, TEAL_MID, SANS)
    p2 = tf.add_paragraph()
    run(p2, title, 20, True, TEAL, SERIF)
    p3 = tf.add_paragraph()
    run(p3, desc, 14, False, GREY, SANS)

tf = tb(s, Inches(0.85), Inches(6.3), Inches(11.4), Inches(0.5))
p = tf.paragraphs[0]
run(p, "Esta diapo es solo el puente. Acá cambiás a la pantalla del sistema.", 13, False, GREY, SANS)

prs.save(str(OUT))
print("OK", OUT)
